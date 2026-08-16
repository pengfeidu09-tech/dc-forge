"""Build a customer-safe enterprise handoff package for the presales solution."""

from __future__ import annotations

from dataclasses import asdict, dataclass, field
from datetime import UTC, datetime
from hashlib import sha256
from html import escape
import json
from pathlib import Path
from typing import Any
from zipfile import ZIP_DEFLATED, ZipFile


PPTX_NAME = "01_DCFORGE_大型汽车制造企业智能招采与采购合规解决方案.pptx"
PDF_NAME = "02_DCFORGE_解决方案建议书.pdf"
HTML_NAME = "03_DCFORGE_解决方案建议书.html"
XLSX_NAME = "04_DCFORGE_实施与验收计划.xlsx"
SECURITY_NAME = "05_DCFORGE_部署安全与集成边界.md"
DEMO_NAME = "06_DCFORGE_演示与操作手册.md"
README_NAME = "README_交付说明.md"
MANIFEST_NAME = "manifest.json"
ARCHIVE_NAME = "DCFORGE_企业智能招采解决方案交付包_v1.0.zip"


@dataclass(frozen=True)
class Scenario:
    name: str
    current_state: str
    ai_assistance: list[str]
    human_boundary: str
    evidence_output: list[str]


@dataclass(frozen=True)
class ImplementationPhase:
    phase: str
    objective: str
    work: list[str]
    outputs: list[str]
    recommended_duration: str
    prerequisite: str


@dataclass(frozen=True)
class AcceptanceMetric:
    dimension: str
    metric: str
    proposed_threshold: str
    evidence: str
    note: str = "建议验收口径，待双方确认"


@dataclass(frozen=True)
class PackageOption:
    name: str
    fit: str
    scope: list[str]
    prerequisite: str
    recommendation: str


@dataclass(frozen=True)
class EnterpriseDeliveryBrief:
    package_version: str
    document_status: str
    client_label: str
    solution_title: str
    prepared_by: str
    known_current_state: list[str]
    pain_points: list[str]
    principles: list[str]
    proposed_goals: list[str]
    workflow: list[str]
    solution_layers: list[dict[str, Any]]
    scenarios: list[Scenario]
    implementation_phases: list[ImplementationPhase]
    acceptance_metrics: list[AcceptanceMetric]
    package_options: list[PackageOption]
    security_boundaries: list[str]
    integration_objects: list[dict[str, str]]
    open_questions: list[str]
    content_boundary: list[str]
    generated_date: str = field(
        default_factory=lambda: datetime.now(UTC).astimezone().date().isoformat()
    )

    def to_dict(self) -> dict[str, Any]:
        return asdict(self)

    @classmethod
    def default_automotive_procurement(cls) -> "EnterpriseDeliveryBrief":
        return cls(
            package_version="1.0",
            document_status="客户交流与试点立项建议稿",
            client_label="某大型汽车制造企业",
            solution_title="大型汽车制造企业智能招采与采购合规解决方案",
            prepared_by="DCForge 项目组",
            known_current_state=[
                "供应商准入主要依赖人工核验材料和制度要求。",
                "询比价过程需要人工收集报价、比对条件并整理说明。",
                "采购合规审查依赖人工查制度、找条款和判断风险。",
                "采购周期较长，关键判断和操作的过程留痕不完整。",
            ],
            pain_points=[
                "资料分散在人员、系统和文档中，重复查找占用采购与合规人员时间。",
                "规则依赖个人经验，遇到特殊采购方式时难以快速形成一致判断。",
                "需求、分析、方案、评审和客户反馈分散在不同入口，版本关系不清楚。",
                "缺少统一证据链，事后复核时难以说明依据、处理人和版本。",
            ],
            principles=[
                "AI 负责材料检查、信息归纳、规则匹配、风险提示和方案草拟。",
                "供应商准入、采购审批、定标和付款仍由企业人员及原有系统完成。",
                "每条建议应能回到客户需求、制度条款、知识来源和人工评审记录。",
                "先在边界清楚的场景做试点，再根据验证结果扩展范围。",
            ],
            proposed_goals=[
                "建立供应商准入、询比价和采购合规审查的统一 AI 助手入口。",
                "把客户沟通转成可确认、可版本化的需求状态，而不是只保留聊天记录。",
                "让企业知识、外部情报和历史案例进入同一套有来源的研究过程。",
                "在内部评审通过后再向客户发布方案，避免内部信息直接对外。",
                "形成可追踪的客户反馈闭环，并保留完整审计记录。",
            ],
            workflow=[
                "客户机会建立",
                "需求分析与确认",
                "外部情报归档",
                "企业知识检索",
                "方案与成果草拟",
                "企业内部评审",
                "客户发布",
                "反馈迭代",
            ],
            solution_layers=[
                {
                    "layer": "客户与员工入口",
                    "components": ["飞书客户机器人", "统一售前工作台", "客户需求与方案中心"],
                    "purpose": "承接客户沟通、内部协作、确认和发布",
                },
                {
                    "layer": "Agent 编排层",
                    "components": ["会话 Agent", "任务编排", "权限与人工审批门禁"],
                    "purpose": "根据用户身份、项目阶段和业务状态选择处理步骤",
                },
                {
                    "layer": "需求与状态层",
                    "components": ["Requirement Intelligence", "业务状态池", "Baseline 与版本差异"],
                    "purpose": "区分客户事实、待确认假设、冲突和缺口",
                },
                {
                    "layer": "知识与工具层",
                    "components": ["企业知识库", "外部情报", "MCP", "业务工具"],
                    "purpose": "提供有权限、有时点、有来源的检索和执行能力",
                },
                {
                    "layer": "方案与交付层",
                    "components": ["三档方案", "内部评审", "客户安全成果", "反馈回流"],
                    "purpose": "把需求和证据转成可审查、可发布、可迭代的成果",
                },
            ],
            scenarios=[
                Scenario(
                    name="供应商准入",
                    current_state="采购人员人工收集和核对资质、证照、制度要求与历史记录。",
                    ai_assistance=[
                        "识别材料类型、有效期、缺失项和字段不一致。",
                        "按准入制度匹配审核项，并给出依据和待补材料。",
                        "形成初审摘要和风险清单，推送给责任人复核。",
                    ],
                    human_boundary="AI 不作最终准入决定。采购负责人按原流程审批。",
                    evidence_output=["材料清单", "规则命中依据", "风险项", "人工处理记录"],
                ),
                Scenario(
                    name="询比价",
                    current_state="采购人员人工整理供应商报价、交付条件、质保和商务条款。",
                    ai_assistance=[
                        "结构化提取价格、税率、交付期、付款和质保条件。",
                        "标出缺项、异常差异和不可直接比较的口径。",
                        "生成比价说明草稿，保留原始报价引用。",
                    ],
                    human_boundary="AI 不替代定标，也不绕过企业采购审批规则。",
                    evidence_output=["报价对比表", "差异说明", "原始报价引用", "人工确认记录"],
                ),
                Scenario(
                    name="采购合规审查",
                    current_state="合规人员人工查找制度、合同和采购材料，逐条判断风险。",
                    ai_assistance=[
                        "识别采购方式、审批链、金额和例外情形。",
                        "检索适用制度条款，提示单一来源、审批缺失等风险。",
                        "生成审查意见草稿并列明证据不足项。",
                    ],
                    human_boundary="高风险事项必须由合规、法务或采购委员会复核。",
                    evidence_output=["适用条款", "风险等级建议", "证据缺口", "最终人工意见"],
                ),
            ],
            implementation_phases=[
                ImplementationPhase(
                    phase="阶段 0：需求基线",
                    objective="把试点范围、规则、数据和验收口径确认清楚。",
                    work=["业务访谈", "制度与样本盘点", "系统边界确认", "RequirementBaseline 评审"],
                    outputs=["确认版需求基线", "数据清单", "接口清单", "试点验收表"],
                    recommended_duration="建议 2 周，待双方排期确认",
                    prerequisite="客户安排采购、合规、IT 和数据安全负责人参加",
                ),
                ImplementationPhase(
                    phase="阶段 1：受控试点",
                    objective="在单一采购范围内验证准入或合规审查闭环。",
                    work=["知识导入", "规则配置", "Agent 联调", "历史样本回放", "用户验收"],
                    outputs=["可演示系统", "试点数据报告", "问题清单", "扩展建议"],
                    recommended_duration="建议 4 周，待数据准备情况确认",
                    prerequisite="提供脱敏样本、制度材料和已标注结果",
                ),
                ImplementationPhase(
                    phase="阶段 2：生产试运行",
                    objective="接入企业系统并在真实流程中受控运行。",
                    work=["SRM、ERP、OA 接口联调", "权限配置", "审计配置", "灰度使用", "运行复盘"],
                    outputs=["生产试运行版本", "运维手册", "审计报告", "阶段验收记录"],
                    recommended_duration="建议 6 至 8 周，待接口和安全评审确认",
                    prerequisite="完成接口、安全、账号和部署环境审批",
                ),
                ImplementationPhase(
                    phase="阶段 3：范围扩展",
                    objective="根据验证结果扩展品类、规则和组织范围。",
                    work=["新增场景", "知识运营", "指标复盘", "模型和规则迭代"],
                    outputs=["扩展版本", "知识运营机制", "季度复盘报告"],
                    recommended_duration="按季度滚动规划",
                    prerequisite="阶段 2 验收通过，并完成扩展范围确认",
                ),
            ],
            acceptance_metrics=[
                AcceptanceMetric(
                    dimension="需求与版本",
                    metric="已确认需求的来源可追溯率",
                    proposed_threshold="建议 100%",
                    evidence="RequirementBaseline、来源引用和版本差异记录",
                ),
                AcceptanceMetric(
                    dimension="供应商准入",
                    metric="约定材料类型的字段提取准确率",
                    proposed_threshold="建议不低于 95%",
                    evidence="双方确认的标注样本和回放报告",
                ),
                AcceptanceMetric(
                    dimension="采购合规",
                    metric="约定重大风险类型的识别召回率",
                    proposed_threshold="建议不低于 90%",
                    evidence="双方确认的风险样本集和评测报告",
                ),
                AcceptanceMetric(
                    dimension="证据与审计",
                    metric="AI 建议的来源引用和操作留痕覆盖率",
                    proposed_threshold="建议 100%",
                    evidence="来源记录、工具调用记录、人工评审和发布记录",
                ),
                AcceptanceMetric(
                    dimension="权限安全",
                    metric="未授权数据泄露事件",
                    proposed_threshold="建议为 0",
                    evidence="权限测试、负向测试和审计日志",
                ),
                AcceptanceMetric(
                    dimension="人工边界",
                    metric="需人工审批事项的人工确认覆盖率",
                    proposed_threshold="建议 100%",
                    evidence="审批记录和人工门禁测试",
                ),
                AcceptanceMetric(
                    dimension="业务效率",
                    metric="约定试点任务的平均人工处理时间变化",
                    proposed_threshold="先采集基线，再由双方确定目标值",
                    evidence="试点前后同口径任务记录",
                ),
            ],
            package_options=[
                PackageOption(
                    name="快速验证",
                    fit="需要先证明可用性和安全边界",
                    scope=["一个优先场景", "离线或脱敏数据", "知识检索", "人工评审", "试点评测"],
                    prerequisite="客户提供样本、制度和人工结果",
                    recommendation="适合作为首次合作入口",
                ),
                PackageOption(
                    name="生产试点",
                    fit="需要进入真实采购流程受控使用",
                    scope=["供应商准入", "询比价", "合规审查", "SRM、ERP、OA 接口", "权限审计"],
                    prerequisite="完成系统接口和数据安全审批",
                    recommendation="本方案建议的目标形态",
                ),
                PackageOption(
                    name="集团扩展",
                    fit="需要跨组织、跨品类运营",
                    scope=["集团知识运营", "多组织权限", "更多采购场景", "指标运营", "持续迭代"],
                    prerequisite="生产试点验收通过",
                    recommendation="根据试点结果分批推进",
                ),
            ],
            security_boundaries=[
                "支持企业私有云或本地部署，具体形态由客户安全评审确定。",
                "采购资料在客户授权的数据域内处理，敏感字段按岗位和场景控制。",
                "知识检索同时检查项目、用户、时间和来源权限。",
                "所有 AI 查询、工具调用、人工修改、评审和发布动作保留审计记录。",
                "对客成果经过内部批准后发布，不显示内部评分、工具标识或系统路径。",
                "AI 不直接执行准入、定标、采购审批、合同签署或付款。",
            ],
            integration_objects=[
                {
                    "system": "SRM",
                    "read": "供应商档案、资质、询价和报价数据",
                    "write": "审查结果、待补材料和风险提示",
                    "mode": "建议通过企业接口或受控数据服务",
                    "boundary": "不绕过 SRM 的业务状态和审批权限",
                },
                {
                    "system": "ERP",
                    "read": "预算、订单、合同、履约和库存相关数据",
                    "write": "原则上不直接写业务凭证",
                    "mode": "只读查询或由客户中台转发",
                    "boundary": "财务和采购凭证仍由原系统生成",
                },
                {
                    "system": "OA",
                    "read": "审批状态、审批意见和组织信息",
                    "write": "AI 建议作为审批附件或参考信息",
                    "mode": "审批接口或消息卡片",
                    "boundary": "审批结论由 OA 和审批人确认",
                },
                {
                    "system": "飞书",
                    "read": "客户与内部人员的授权会话",
                    "write": "客户回复、工作台入口和待办通知",
                    "mode": "应用机器人长连接",
                    "boundary": "客户对话与内部工作台按身份隔离",
                },
            ],
            open_questions=[
                "首个试点优先选择供应商准入、询比价还是采购合规审查？",
                "试点涉及哪些采购品类、组织和用户角色？",
                "现有审批规则、金额阈值和例外情形以哪套制度为准？",
                "可用于试点的历史材料、制度、报价和标注结果有哪些？",
                "SRM、ERP、OA 可以提供哪些接口和测试环境？",
                "部署、网络、模型、日志、脱敏和数据留存要求是什么？",
                "双方采用哪些样本集、计算口径和验收阈值？",
            ],
            content_boundary=[
                "本材料基于前期场景描述整理，尚未替代双方签署的需求确认文件。",
                "材料中的指标是建议验收口径，不是当前客户已经取得的业务成果。",
                "实施周期是建议排期，需根据数据、接口、安全评审和人员安排确认。",
                "正式范围、价格、服务等级和责任边界以双方合同及附件为准。",
            ],
        )


@dataclass(frozen=True)
class DeliveryBuildResult:
    output_dir: Path
    archive_path: Path
    manifest: dict[str, Any]


class EnterpriseDeliveryPackageBuilder:
    """Create the editable and printable files used in an enterprise handoff."""

    def __init__(self, brief: EnterpriseDeliveryBrief) -> None:
        self.brief = brief

    def build(self, output_dir: Path | str) -> DeliveryBuildResult:
        root = Path(output_dir)
        root.mkdir(parents=True, exist_ok=True)
        paths = {
            PPTX_NAME: root / PPTX_NAME,
            PDF_NAME: root / PDF_NAME,
            HTML_NAME: root / HTML_NAME,
            XLSX_NAME: root / XLSX_NAME,
            SECURITY_NAME: root / SECURITY_NAME,
            DEMO_NAME: root / DEMO_NAME,
            README_NAME: root / README_NAME,
        }
        self._write_pptx(paths[PPTX_NAME])
        self._write_pdf(paths[PDF_NAME])
        paths[HTML_NAME].write_text(self._html(), encoding="utf-8")
        self._write_xlsx(paths[XLSX_NAME])
        paths[SECURITY_NAME].write_text(self._security_markdown(), encoding="utf-8")
        paths[DEMO_NAME].write_text(self._demo_markdown(), encoding="utf-8")
        paths[README_NAME].write_text(self._readme_markdown(), encoding="utf-8")

        manifest = self._manifest(paths)
        manifest_path = root / MANIFEST_NAME
        manifest_path.write_text(
            json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
        archive_path = root.parent / ARCHIVE_NAME
        with ZipFile(archive_path, "w", ZIP_DEFLATED) as archive:
            for name in [*paths, MANIFEST_NAME]:
                archive.write(root / name, name)
        return DeliveryBuildResult(root, archive_path, manifest)

    def _manifest(self, paths: dict[str, Path]) -> dict[str, Any]:
        files = []
        for name, path in paths.items():
            payload = path.read_bytes()
            files.append(
                {
                    "name": name,
                    "size_bytes": len(payload),
                    "sha256": sha256(payload).hexdigest(),
                }
            )
        return {
            "package_name": "DCForge 企业智能招采与采购合规解决方案交付包",
            "package_version": self.brief.package_version,
            "document_status": self.brief.document_status,
            "generated_at": datetime.now(UTC).isoformat(),
            "business_results_status": "未验证",
            "intended_use": ["客户正式汇报", "技术交流", "试点立项", "验收口径讨论"],
            "content_boundary": self.brief.content_boundary,
            "files": files,
        }

    @staticmethod
    def _require(module_name: str, installation: str) -> Any:
        try:
            return __import__(module_name)
        except ImportError as error:
            raise RuntimeError(
                f"Enterprise delivery generation requires {installation}"
            ) from error

    def _write_pptx(self, path: Path) -> None:
        self._require("pptx", "python-pptx")
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt

        prs = Presentation()
        prs.slide_width = Inches(13.333333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        navy = RGBColor(11, 23, 57)
        blue = RGBColor(30, 94, 255)
        cyan = RGBColor(24, 183, 201)
        green = RGBColor(16, 185, 129)
        orange = RGBColor(245, 158, 11)
        red = RGBColor(220, 38, 38)
        ink = RGBColor(23, 32, 51)
        muted = RGBColor(93, 108, 133)
        line = RGBColor(219, 227, 240)
        light = RGBColor(244, 247, 251)
        white = RGBColor(255, 255, 255)

        def rect(slide, x, y, w, h, fill, radius=True, border=None):
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill
            shape.line.color.rgb = border or fill
            return shape

        def textbox(
            slide,
            text,
            x,
            y,
            w,
            h,
            *,
            size=18,
            color=ink,
            bold=False,
            align=PP_ALIGN.LEFT,
            valign=MSO_ANCHOR.TOP,
            font="PingFang SC",
            margin=0.04,
        ):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            frame = box.text_frame
            frame.clear()
            frame.margin_left = frame.margin_right = Inches(margin)
            frame.margin_top = frame.margin_bottom = Inches(margin)
            frame.vertical_anchor = valign
            paragraph = frame.paragraphs[0]
            paragraph.text = text
            paragraph.alignment = align
            paragraph.font.name = font
            paragraph.font.size = Pt(size)
            paragraph.font.bold = bold
            paragraph.font.color.rgb = color
            return box

        def bullets(slide, values, x, y, w, h, *, size=16, color=ink, gap=7):
            box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
            frame = box.text_frame
            frame.clear()
            frame.margin_left = frame.margin_right = Inches(0.03)
            frame.margin_top = frame.margin_bottom = Inches(0.02)
            for index, value in enumerate(values):
                paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                paragraph.text = f"• {value}"
                paragraph.font.name = "PingFang SC"
                paragraph.font.size = Pt(size)
                paragraph.font.color.rgb = color
                paragraph.space_after = Pt(gap)
            return box

        def base_slide(title, kicker=None, dark=False):
            slide = prs.slides.add_slide(blank)
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = navy if dark else light
            if dark:
                rect(slide, 0, 0, 0.14, 7.5, cyan, radius=False)
            else:
                rect(slide, 0, 0, 13.333, 0.1, blue, radius=False)
            if kicker:
                textbox(slide, kicker, 0.65, 0.34, 3.0, 0.3, size=10, color=cyan if dark else blue, bold=True)
            textbox(slide, title, 0.65, 0.68 if kicker else 0.42, 11.8, 0.55, size=26, color=white if dark else navy, bold=True)
            textbox(
                slide,
                f"DCForge · {self.brief.document_status} · v{self.brief.package_version}",
                0.65,
                7.15,
                8.4,
                0.18,
                size=8,
                color=RGBColor(170, 184, 210) if dark else muted,
            )
            textbox(slide, str(len(prs.slides)), 12.2, 7.12, 0.45, 0.2, size=8, color=RGBColor(170, 184, 210) if dark else muted, align=PP_ALIGN.RIGHT)
            return slide

        def card(slide, title, body, x, y, w, h, accent=blue, number=None):
            rect(slide, x, y, w, h, white, border=line)
            rect(slide, x, y, 0.08, h, accent, radius=False)
            if number is not None:
                rect(slide, x + 0.25, y + 0.22, 0.46, 0.46, accent)
                textbox(slide, str(number), x + 0.25, y + 0.28, 0.46, 0.25, size=12, color=white, bold=True, align=PP_ALIGN.CENTER)
                title_x = x + 0.84
                title_w = w - 1.05
            else:
                title_x = x + 0.28
                title_w = w - 0.5
            textbox(slide, title, title_x, y + 0.22, title_w, 0.38, size=16, color=navy, bold=True)
            textbox(slide, body, x + 0.28, y + 0.78, w - 0.52, h - 0.98, size=12, color=muted)

        slide = base_slide("", dark=True)
        textbox(slide, "DCFORGE", 0.85, 0.8, 3.0, 0.45, size=14, color=cyan, bold=True)
        textbox(slide, "大型汽车制造企业\n智能招采与采购合规解决方案", 0.85, 1.58, 8.45, 1.55, size=29, color=white, bold=True)
        textbox(slide, "从客户沟通到需求基线、企业知识、方案评审和客户发布的完整售前闭环", 0.88, 3.32, 8.45, 0.74, size=17, color=RGBColor(205, 218, 242))
        rect(slide, 0.88, 4.45, 4.8, 0.68, RGBColor(23, 46, 96))
        textbox(slide, self.brief.document_status, 1.12, 4.66, 4.3, 0.25, size=12, color=white, bold=True)
        textbox(slide, f"面向：{self.brief.client_label}", 0.88, 5.55, 5.7, 0.35, size=13, color=RGBColor(205, 218, 242))
        textbox(slide, f"{self.brief.prepared_by} · {self.brief.generated_date}", 0.88, 6.02, 5.7, 0.35, size=12, color=RGBColor(160, 180, 215))
        rect(slide, 9.9, 1.75, 2.35, 3.8, RGBColor(17, 38, 82))
        for idx, label in enumerate(["客户", "Agent", "需求", "知识", "方案", "评审", "发布"]):
            color = cyan if idx in {0, 6} else blue
            rect(slide, 10.3, 2.05 + idx * 0.45, 1.55, 0.3, color)
            textbox(slide, label, 10.3, 2.12 + idx * 0.45, 1.55, 0.15, size=8, color=white, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("我们目前如何理解贵司", "01 客户现状")
        textbox(slide, "以下内容来自前期场景描述，需在项目启动阶段形成双方确认的 RequirementBaseline。", 0.68, 1.35, 11.9, 0.4, size=12, color=orange, bold=True)
        positions = [(0.68, 1.95), (6.82, 1.95), (0.68, 4.25), (6.82, 4.25)]
        for idx, (value, (x, y)) in enumerate(zip(self.brief.known_current_state, positions), 1):
            card(slide, ["供应商准入", "询比价", "合规审查", "采购运营"][idx - 1], value, x, y, 5.82, 1.75, [blue, cyan, orange, green][idx - 1], idx)

        slide = base_slide("本次建议聚焦的问题", "02 问题定义")
        for idx, problem in enumerate(self.brief.pain_points, 1):
            y = 1.45 + (idx - 1) * 1.28
            rect(slide, 0.72, y, 0.6, 0.6, [blue, cyan, orange, red][idx - 1])
            textbox(slide, f"0{idx}", 0.72, y + 0.17, 0.6, 0.2, size=11, color=white, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, problem, 1.58, y + 0.06, 10.7, 0.55, size=16, color=ink, bold=True)
        textbox(slide, "问题不在于缺少一个聊天机器人，而在于客户事实、知识证据、内部判断和对客成果没有形成同一条可追踪链路。", 0.75, 6.55, 11.8, 0.38, size=13, color=navy, bold=True)

        slide = base_slide("建设目标和验收口径说明", "03 建设目标")
        for idx, goal in enumerate(self.brief.proposed_goals, 1):
            x = 0.68 + ((idx - 1) % 3) * 4.18
            y = 1.45 + ((idx - 1) // 3) * 2.35
            card(slide, f"目标 {idx}", goal, x, y, 3.82, 1.85, [blue, cyan, green, orange, blue][idx - 1])
        rect(slide, 4.86, 4.0, 3.82, 1.85, navy)
        textbox(slide, "口径说明", 5.18, 4.3, 3.15, 0.35, size=17, color=white, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, "本材料中的目标值是建议验收口径。双方应先确认基线、样本和计算方法，再判断是否达成。", 5.15, 4.88, 3.24, 0.64, size=12, color=RGBColor(213, 225, 247), align=PP_ALIGN.CENTER)

        slide = base_slide("方案总览", "04 方案架构")
        layer_colors = [cyan, blue, RGBColor(83, 103, 232), orange, green]
        for idx, layer in enumerate(self.brief.solution_layers):
            y = 1.38 + idx * 1.05
            rect(slide, 0.7, y, 2.15, 0.72, layer_colors[idx])
            textbox(slide, layer["layer"], 0.86, y + 0.2, 1.83, 0.25, size=13, color=white, bold=True, align=PP_ALIGN.CENTER)
            rect(slide, 3.15, y, 6.1, 0.72, white, border=line)
            textbox(slide, "  ·  ".join(layer["components"]), 3.38, y + 0.17, 5.7, 0.3, size=13, color=navy, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, layer["purpose"], 9.55, y + 0.13, 2.95, 0.42, size=11, color=muted)
        textbox(slide, "核心原则：客户信息先进入需求和状态管理，再调用知识与工具。成果经过内部评审后才对外发布。", 0.72, 6.75, 11.8, 0.28, size=12, color=blue, bold=True)

        slide = base_slide("从客户机会到反馈迭代的完整闭环", "05 业务流程", dark=True)
        for idx, label in enumerate(self.brief.workflow, 1):
            col = (idx - 1) % 4
            row = (idx - 1) // 4
            x = 0.78 + col * 3.1
            y = 1.55 + row * 2.25
            rect(slide, x, y, 2.62, 1.35, RGBColor(19, 42, 88), border=RGBColor(48, 81, 146))
            rect(slide, x + 0.18, y + 0.18, 0.48, 0.48, cyan if idx in {1, 8} else blue)
            textbox(slide, str(idx), x + 0.18, y + 0.32, 0.48, 0.16, size=10, color=white, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, label, x + 0.78, y + 0.22, 1.62, 0.55, size=14, color=white, bold=True)
            textbox(slide, "事实与版本留痕", x + 0.78, y + 0.84, 1.6, 0.25, size=9, color=RGBColor(177, 196, 231))
        textbox(slide, "每一阶段都由持久化事实推导状态，前端不能手工把项目标成完成。", 0.8, 6.3, 11.7, 0.35, size=13, color=cyan, bold=True, align=PP_ALIGN.CENTER)

        for slide_number, scenario in enumerate(self.brief.scenarios, 6):
            slide = base_slide(scenario.name, f"0{slide_number} 重点场景")
            card(slide, "当前做法", scenario.current_state, 0.7, 1.4, 3.55, 1.7, orange)
            rect(slide, 4.63, 1.4, 4.2, 3.95, white, border=line)
            textbox(slide, "AI 可以承担的工作", 4.95, 1.73, 3.55, 0.35, size=17, color=navy, bold=True)
            bullets(slide, scenario.ai_assistance, 4.95, 2.35, 3.5, 2.35, size=13, color=ink, gap=11)
            card(slide, "人工决策边界", scenario.human_boundary, 9.2, 1.4, 3.4, 1.7, red)
            card(slide, "系统交付证据", "、".join(scenario.evidence_output), 0.7, 3.55, 3.55, 1.8, green)
            rect(slide, 9.2, 3.55, 3.4, 1.8, navy)
            textbox(slide, "输出原则", 9.5, 3.86, 2.78, 0.32, size=16, color=white, bold=True)
            textbox(slide, "给出建议的同时，显示依据、缺口、版本和处理人。证据不足时明确提示人工补充。", 9.5, 4.42, 2.75, 0.62, size=12, color=RGBColor(213, 225, 247))
            textbox(slide, "试点从一类材料、一套规则和一条人工审批链开始，具体范围待双方确认。", 0.72, 6.45, 11.8, 0.38, size=13, color=blue, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("Agent、知识库、MCP 与业务工具", "09 技术架构")
        columns = [
            ("入口", ["飞书客户机器人", "内部统一工作台", "客户方案中心"], cyan),
            ("Agent", ["身份识别", "项目状态判断", "任务编排", "人工门禁"], blue),
            ("智能能力", ["Requirement Intelligence", "企业知识检索", "方案生成", "成果编辑"], RGBColor(83, 103, 232)),
            ("MCP 与工具", ["项目查询", "来源检索", "供应商分析", "文档审查", "方案编译"], orange),
            ("治理", ["权限", "时间边界", "来源引用", "审计日志", "内部评审"], green),
        ]
        for idx, (title, items, accent) in enumerate(columns):
            x = 0.52 + idx * 2.56
            rect(slide, x, 1.48, 2.22, 4.65, white, border=line)
            rect(slide, x, 1.48, 2.22, 0.64, accent)
            textbox(slide, title, x + 0.16, 1.68, 1.9, 0.25, size=14, color=white, bold=True, align=PP_ALIGN.CENTER)
            for row, item in enumerate(items):
                rect(slide, x + 0.18, 2.38 + row * 0.66, 1.86, 0.46, light, border=line)
                textbox(slide, item, x + 0.26, 2.52 + row * 0.66, 1.7, 0.17, size=10, color=ink, align=PP_ALIGN.CENTER)
        textbox(slide, "Agent 不直接读取所有数据。每次检索和工具调用都需要项目、用户、时间和权限上下文。", 0.7, 6.55, 11.9, 0.35, size=13, color=navy, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("业务状态池和需求版本机制", "10 需求与状态")
        stages = [
            ("客户原话", "保留会话与材料来源", cyan),
            ("需求识别", "提取事实、假设、冲突和缺口", blue),
            ("客户确认", "形成正式 RequirementBaseline", green),
            ("方案编译", "绑定 Baseline、研究和方案版本", orange),
            ("反馈迭代", "新信息生成新 State 与差异", RGBColor(83, 103, 232)),
        ]
        for idx, (title, body, accent) in enumerate(stages):
            x = 0.65 + idx * 2.5
            rect(slide, x, 1.62, 2.05, 1.55, white, border=line)
            rect(slide, x, 1.62, 2.05, 0.12, accent, radius=False)
            textbox(slide, title, x + 0.18, 1.95, 1.7, 0.33, size=15, color=navy, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, body, x + 0.2, 2.48, 1.65, 0.42, size=10, color=muted, align=PP_ALIGN.CENTER)
            if idx < len(stages) - 1:
                textbox(slide, "→", x + 2.12, 2.17, 0.34, 0.35, size=18, color=blue, bold=True, align=PP_ALIGN.CENTER)
        card(slide, "状态不等于客户事实", "外部情报、企业知识和 Agent 推断只作为证据或待确认假设，不能自动写入客户确认基线。", 0.72, 3.82, 3.75, 1.75, orange)
        card(slide, "编辑会使旧批准失效", "成果稿修改后修订号增加。发布前必须对当前修订重新批准。", 4.8, 3.82, 3.75, 1.75, red)
        card(slide, "客户反馈回到需求分析", "客户补充和纠正形成新的需求版本，内部员工可以查看变化和处理记录。", 8.88, 3.82, 3.75, 1.75, green)

        slide = base_slide("企业员工和客户的双入口", "11 使用入口")
        rect(slide, 0.7, 1.42, 5.75, 4.9, white, border=line)
        rect(slide, 6.88, 1.42, 5.75, 4.9, white, border=line)
        rect(slide, 0.7, 1.42, 5.75, 0.75, blue)
        rect(slide, 6.88, 1.42, 5.75, 0.75, cyan)
        textbox(slide, "企业内部统一售前工作台", 1.0, 1.66, 5.15, 0.28, size=17, color=white, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, "客户需求与方案中心", 7.18, 1.66, 5.15, 0.28, size=17, color=white, bold=True, align=PP_ALIGN.CENTER)
        bullets(slide, ["飞书沟通时间线", "RequirementState、缺口和 Baseline", "资料、外部情报与知识研究", "三档方案和成果稿编辑", "内部批准、驳回和发布记录"], 1.08, 2.62, 4.95, 2.95, size=14, gap=10)
        bullets(slide, ["查看当前需求理解", "确认、否定或补充需求", "查看正式发布的解决方案", "下载可编辑 HTML 成果", "提交反馈并进入新一轮分析"], 7.25, 2.62, 4.95, 2.95, size=14, gap=10)
        textbox(slide, "客户页面不显示内部评分、内部知识路径、工具标识和权限配置。", 6.98, 5.82, 5.5, 0.3, size=11, color=red, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("数据、安全与人工决策边界", "12 安全治理")
        for idx, boundary in enumerate(self.brief.security_boundaries, 1):
            col = (idx - 1) % 2
            row = (idx - 1) // 2
            x = 0.7 + col * 6.15
            y = 1.38 + row * 1.62
            card(slide, ["部署与数据域", "岗位权限", "检索权限", "审计记录", "对客发布", "人工审批"][idx - 1], boundary, x, y, 5.75, 1.26, [blue, cyan, green, orange, RGBColor(83, 103, 232), red][idx - 1], idx)

        slide = base_slide("分阶段实施路线", "13 实施计划", dark=True)
        for idx, phase in enumerate(self.brief.implementation_phases):
            x = 0.66 + idx * 3.14
            rect(slide, x, 1.48, 2.78, 4.86, RGBColor(19, 42, 88), border=RGBColor(48, 81, 146))
            rect(slide, x, 1.48, 2.78, 0.68, [cyan, blue, orange, green][idx])
            textbox(slide, phase.phase, x + 0.14, 1.7, 2.5, 0.24, size=13, color=white, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, phase.objective, x + 0.24, 2.45, 2.3, 0.65, size=12, color=white, bold=True, align=PP_ALIGN.CENTER)
            bullets(slide, phase.work, x + 0.24, 3.35, 2.3, 1.62, size=10, color=RGBColor(205, 218, 242), gap=5)
            textbox(slide, phase.recommended_duration, x + 0.24, 5.4, 2.3, 0.48, size=9, color=cyan, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, "建议排期用于资源规划，正式周期取决于数据、接口、安全评审和双方人员安排。", 0.8, 6.64, 11.7, 0.3, size=12, color=cyan, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("建议验收指标", "14 验收口径")
        headers = ["维度", "指标", "建议口径", "验收证据"]
        widths = [1.45, 3.3, 2.45, 4.5]
        x_positions = [0.68]
        for width in widths[:-1]:
            x_positions.append(x_positions[-1] + width)
        for x, width, header in zip(x_positions, widths, headers):
            rect(slide, x, 1.35, width, 0.52, navy, radius=False)
            textbox(slide, header, x + 0.06, 1.52, width - 0.12, 0.18, size=10, color=white, bold=True, align=PP_ALIGN.CENTER)
        for row, metric in enumerate(self.brief.acceptance_metrics):
            y = 1.87 + row * 0.68
            values = [metric.dimension, metric.metric, metric.proposed_threshold, metric.evidence]
            for x, width, value in zip(x_positions, widths, values):
                rect(slide, x, y, width, 0.68, white if row % 2 == 0 else RGBColor(237, 242, 250), radius=False, border=line)
                textbox(slide, value, x + 0.08, y + 0.12, width - 0.16, 0.4, size=8.5, color=ink, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        textbox(slide, "所有阈值均待双方确认。验收前应冻结样本集、计算口径、人工金标准和异常处理规则。", 0.7, 6.72, 11.9, 0.25, size=11, color=orange, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("三档建设范围", "15 建设选择")
        for idx, option in enumerate(self.brief.package_options):
            x = 0.68 + idx * 4.18
            accent = [cyan, blue, green][idx]
            rect(slide, x, 1.42, 3.76, 4.95, white, border=accent if idx == 1 else line)
            rect(slide, x, 1.42, 3.76, 0.74, accent)
            textbox(slide, option.name, x + 0.2, 1.66, 3.36, 0.26, size=17, color=white, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, option.fit, x + 0.28, 2.45, 3.2, 0.6, size=12, color=navy, bold=True, align=PP_ALIGN.CENTER)
            bullets(slide, option.scope, x + 0.35, 3.28, 3.0, 1.72, size=11, gap=6)
            textbox(slide, option.prerequisite, x + 0.3, 5.3, 3.15, 0.4, size=9, color=muted, align=PP_ALIGN.CENTER)
            textbox(slide, option.recommendation, x + 0.3, 5.84, 3.15, 0.25, size=10, color=accent, bold=True, align=PP_ALIGN.CENTER)
        textbox(slide, "建议先完成需求基线和受控试点，再确认生产范围与商务报价。", 0.7, 6.68, 11.9, 0.3, size=12, color=navy, bold=True, align=PP_ALIGN.CENTER)

        slide = base_slide("项目风险与双方待确认事项", "16 风险与前提")
        rect(slide, 0.68, 1.38, 5.8, 5.25, white, border=line)
        rect(slide, 6.85, 1.38, 5.8, 5.25, white, border=line)
        textbox(slide, "主要项目风险", 1.0, 1.74, 5.15, 0.35, size=18, color=navy, bold=True)
        bullets(slide, ["制度和样本不完整会影响识别与评测结果。", "接口或权限审批延迟会影响生产试运行排期。", "不同组织和品类规则差异较大，需要分范围配置。", "没有冻结样本和计算口径时，效率与准确率无法客观验收。", "AI 建议必须保留人工复核，避免越过采购治理边界。"], 1.0, 2.38, 5.05, 3.45, size=13, gap=10)
        textbox(slide, "需要双方确认", 7.18, 1.74, 5.15, 0.35, size=18, color=navy, bold=True)
        bullets(slide, self.brief.open_questions[:6], 7.18, 2.38, 5.05, 3.65, size=12, gap=8)

        slide = base_slide("建议下一步", "17 行动安排", dark=True)
        next_steps = [
            ("01", "场景确认会", "确定首个试点场景、参与部门和负责人"),
            ("02", "资料与系统盘点", "确认制度、样本、接口、部署和权限边界"),
            ("03", "需求基线评审", "冻结试点范围、样本集、指标和人工决策点"),
            ("04", "受控试点启动", "按确认计划完成配置、联调、回放和验收"),
        ]
        for idx, (number, title, body) in enumerate(next_steps):
            y = 1.35 + idx * 1.28
            rect(slide, 0.88, y, 0.72, 0.72, cyan if idx in {0, 3} else blue)
            textbox(slide, number, 0.88, y + 0.22, 0.72, 0.22, size=11, color=white, bold=True, align=PP_ALIGN.CENTER)
            textbox(slide, title, 1.93, y + 0.03, 2.35, 0.33, size=18, color=white, bold=True)
            textbox(slide, body, 4.55, y + 0.05, 7.2, 0.4, size=14, color=RGBColor(205, 218, 242))
        rect(slide, 0.88, 6.55, 11.3, 0.07, cyan, radius=False)
        textbox(slide, "本材料可直接用于首次正式汇报和试点立项讨论。正式实施以双方确认的需求基线、合同和验收附件为准。", 0.9, 6.78, 11.25, 0.35, size=12, color=cyan, bold=True, align=PP_ALIGN.CENTER)

        prs.core_properties.title = self.brief.solution_title
        prs.core_properties.subject = self.brief.document_status
        prs.core_properties.author = self.brief.prepared_by
        prs.core_properties.comments = "指标为建议验收口径，业务成果尚未验证。"
        prs.save(path)

    def _write_pdf(self, path: Path) -> None:
        self._require("reportlab", "reportlab")
        from reportlab.lib import colors
        from reportlab.lib.enums import TA_CENTER
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
        from reportlab.lib.units import mm
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        from reportlab.platypus import (
            PageBreak,
            Paragraph,
            SimpleDocTemplate,
            Spacer,
            Table,
            TableStyle,
        )

        pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
        styles = getSampleStyleSheet()
        title = ParagraphStyle(
            "CnTitle", parent=styles["Title"], fontName="STSong-Light", fontSize=25,
            leading=34, textColor=colors.HexColor("#0B1739"), alignment=TA_CENTER,
            spaceAfter=12,
        )
        subtitle = ParagraphStyle(
            "CnSubtitle", parent=styles["Normal"], fontName="STSong-Light", fontSize=12,
            leading=20, textColor=colors.HexColor("#5D6C85"), alignment=TA_CENTER,
        )
        h1 = ParagraphStyle(
            "CnH1", parent=styles["Heading1"], fontName="STSong-Light", fontSize=18,
            leading=25, textColor=colors.HexColor("#0B1739"), spaceBefore=10, spaceAfter=8,
        )
        h2 = ParagraphStyle(
            "CnH2", parent=styles["Heading2"], fontName="STSong-Light", fontSize=13,
            leading=19, textColor=colors.HexColor("#1E5EFF"), spaceBefore=8, spaceAfter=5,
        )
        body = ParagraphStyle(
            "CnBody", parent=styles["BodyText"], fontName="STSong-Light", fontSize=10,
            leading=17, textColor=colors.HexColor("#172033"), spaceAfter=5,
        )
        small = ParagraphStyle(
            "CnSmall", parent=body, fontSize=8.5, leading=13,
            textColor=colors.HexColor("#5D6C85"),
        )
        notice = ParagraphStyle(
            "CnNotice", parent=body, fontSize=10.5, leading=17,
            textColor=colors.HexColor("#8A4B08"), borderColor=colors.HexColor("#F2C76E"),
            borderWidth=0.8, borderPadding=8, backColor=colors.HexColor("#FFF8E7"),
        )

        doc = SimpleDocTemplate(
            str(path), pagesize=A4, rightMargin=18 * mm, leftMargin=18 * mm,
            topMargin=17 * mm, bottomMargin=16 * mm,
            title=self.brief.solution_title, author=self.brief.prepared_by,
        )
        story: list[Any] = [
            Spacer(1, 32 * mm),
            Paragraph(self.brief.solution_title, title),
            Spacer(1, 8 * mm),
            Paragraph("从客户沟通到需求基线、企业知识、方案评审和客户发布的完整售前闭环", subtitle),
            Spacer(1, 18 * mm),
            Paragraph(self.brief.document_status, notice),
            Spacer(1, 18 * mm),
            Paragraph(f"面向：{self.brief.client_label}", subtitle),
            Paragraph(f"编制：{self.brief.prepared_by}", subtitle),
            Paragraph(f"日期：{self.brief.generated_date}　版本：{self.brief.package_version}", subtitle),
            PageBreak(),
        ]

        def add_bullets(values: list[str]) -> None:
            for value in values:
                story.append(Paragraph(f"• {escape(value)}", body))

        story.extend([
            Paragraph("一、材料说明", h1),
            Paragraph("本建议书用于客户正式汇报、技术交流、试点立项和验收口径讨论。项目启动后，双方应通过 RequirementBaseline 确认正式需求。", body),
            Paragraph("材料中的目标值是建议验收口径，不是当前客户已经取得的业务成果。实施范围、周期、价格和服务等级以双方后续确认文件为准。", notice),
            Paragraph("二、客户现状理解", h1),
        ])
        add_bullets(self.brief.known_current_state)
        story.append(Paragraph("三、主要问题", h1))
        add_bullets(self.brief.pain_points)
        story.append(Paragraph("四、建设原则", h1))
        add_bullets(self.brief.principles)
        story.append(PageBreak())
        story.append(Paragraph("五、总体方案", h1))
        layer_data = [[Paragraph("层级", small), Paragraph("组成", small), Paragraph("作用", small)]]
        for layer in self.brief.solution_layers:
            layer_data.append([
                Paragraph(escape(layer["layer"]), small),
                Paragraph(escape("、".join(layer["components"])), small),
                Paragraph(escape(layer["purpose"]), small),
            ])
        layer_table = Table(layer_data, colWidths=[31 * mm, 78 * mm, 55 * mm], repeatRows=1)
        layer_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0B1739")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor("#D6DFEE")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F7FB")]),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
            ("RIGHTPADDING", (0, 0), (-1, -1), 6),
            ("TOPPADDING", (0, 0), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
        ]))
        story.extend([layer_table, Spacer(1, 6 * mm), Paragraph("业务闭环", h2)])
        add_bullets([f"{index + 1}. {value}" for index, value in enumerate(self.brief.workflow)])

        story.append(PageBreak())
        story.append(Paragraph("六、重点业务场景", h1))
        for scenario in self.brief.scenarios:
            story.append(Paragraph(scenario.name, h2))
            story.append(Paragraph(f"当前做法：{escape(scenario.current_state)}", body))
            add_bullets(scenario.ai_assistance)
            story.append(Paragraph(f"人工边界：{escape(scenario.human_boundary)}", notice))
            story.append(Paragraph(f"交付证据：{escape('、'.join(scenario.evidence_output))}", body))

        story.append(PageBreak())
        story.append(Paragraph("七、数据、安全和人工决策边界", h1))
        add_bullets(self.brief.security_boundaries)
        story.append(Paragraph("系统集成建议", h2))
        integration_data = [[Paragraph(value, small) for value in ["系统", "读取", "回写", "方式", "边界"]]]
        for item in self.brief.integration_objects:
            integration_data.append([
                Paragraph(escape(item[key]), small)
                for key in ["system", "read", "write", "mode", "boundary"]
            ])
        integration_table = Table(integration_data, colWidths=[19 * mm, 37 * mm, 36 * mm, 37 * mm, 40 * mm], repeatRows=1)
        integration_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1E5EFF")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D6DFEE")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F7FB")]),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.append(integration_table)

        story.append(PageBreak())
        story.append(Paragraph("八、实施路线", h1))
        for phase in self.brief.implementation_phases:
            story.append(Paragraph(phase.phase, h2))
            story.append(Paragraph(phase.objective, body))
            add_bullets(phase.work)
            story.append(Paragraph(f"阶段输出：{escape('、'.join(phase.outputs))}", body))
            story.append(Paragraph(f"建议周期：{escape(phase.recommended_duration)}", small))
            story.append(Paragraph(f"前提：{escape(phase.prerequisite)}", small))

        story.append(PageBreak())
        story.append(Paragraph("九、建议验收指标", h1))
        metric_data = [[Paragraph(value, small) for value in ["维度", "指标", "建议口径", "证据"]]]
        for metric in self.brief.acceptance_metrics:
            metric_data.append([
                Paragraph(escape(metric.dimension), small),
                Paragraph(escape(metric.metric), small),
                Paragraph(escape(metric.proposed_threshold), small),
                Paragraph(escape(metric.evidence), small),
            ])
        metric_table = Table(metric_data, colWidths=[27 * mm, 53 * mm, 38 * mm, 51 * mm], repeatRows=1)
        metric_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0B1739")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#D6DFEE")),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#F4F7FB")]),
            ("LEFTPADDING", (0, 0), (-1, -1), 5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 5),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
        ]))
        story.extend([metric_table, Spacer(1, 5 * mm), Paragraph("所有阈值均待双方确认。验收前应冻结样本集、计算口径、人工金标准和异常处理规则。", notice)])

        story.append(PageBreak())
        story.append(Paragraph("十、建设范围选择", h1))
        for option in self.brief.package_options:
            story.append(Paragraph(option.name, h2))
            story.append(Paragraph(f"适用情况：{escape(option.fit)}", body))
            add_bullets(option.scope)
            story.append(Paragraph(f"前提：{escape(option.prerequisite)}", small))
            story.append(Paragraph(option.recommendation, notice))
        story.append(Paragraph("十一、双方待确认事项", h1))
        add_bullets(self.brief.open_questions)
        story.append(Paragraph("十二、内容边界", h1))
        add_bullets(self.brief.content_boundary)
        doc.build(story)

    def _write_xlsx(self, path: Path) -> None:
        self._require("openpyxl", "openpyxl")
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Border, Font, PatternFill, Side

        workbook = Workbook()
        workbook.remove(workbook.active)
        navy = "0B1739"
        blue = "1E5EFF"
        light = "F4F7FB"
        white = "FFFFFF"
        border = Border(
            left=Side(style="thin", color="D6DFEE"),
            right=Side(style="thin", color="D6DFEE"),
            top=Side(style="thin", color="D6DFEE"),
            bottom=Side(style="thin", color="D6DFEE"),
        )

        def configure(sheet, widths: list[int]) -> None:
            sheet.freeze_panes = "A2"
            sheet.auto_filter.ref = sheet.dimensions
            for cell in sheet[1]:
                cell.fill = PatternFill("solid", fgColor=navy)
                cell.font = Font(color=white, bold=True)
                cell.alignment = Alignment(horizontal="center", vertical="center")
            for row in sheet.iter_rows():
                for cell in row:
                    cell.border = border
                    cell.alignment = Alignment(vertical="top", wrap_text=True)
                    if cell.row > 1 and cell.row % 2 == 0:
                        cell.fill = PatternFill("solid", fgColor=light)
            for index, width in enumerate(widths, 1):
                sheet.column_dimensions[chr(64 + index)].width = width
            sheet.row_dimensions[1].height = 26

        plan = workbook.create_sheet("实施计划")
        plan.append(["阶段", "目标", "主要工作", "阶段输出", "建议周期", "前提条件", "状态"])
        for phase in self.brief.implementation_phases:
            plan.append([
                phase.phase,
                phase.objective,
                "\n".join(f"• {value}" for value in phase.work),
                "\n".join(f"• {value}" for value in phase.outputs),
                phase.recommended_duration,
                phase.prerequisite,
                "待双方确认",
            ])
        configure(plan, [22, 34, 42, 36, 27, 38, 18])

        metrics = workbook.create_sheet("验收指标")
        metrics.append(["维度", "指标", "建议阈值", "验收证据", "状态", "备注"])
        for metric in self.brief.acceptance_metrics:
            metrics.append([
                metric.dimension,
                metric.metric,
                metric.proposed_threshold,
                metric.evidence,
                "未验证",
                metric.note,
            ])
        configure(metrics, [20, 38, 30, 46, 16, 30])

        interfaces = workbook.create_sheet("接口清单")
        interfaces.append(["系统", "读取数据", "建议回写", "接入方式", "安全与业务边界", "接口负责人", "状态"])
        for item in self.brief.integration_objects:
            interfaces.append([
                item["system"], item["read"], item["write"], item["mode"], item["boundary"], "待客户指定", "待确认"
            ])
        configure(interfaces, [16, 38, 38, 34, 45, 20, 16])

        actions = workbook.create_sheet("双方待办")
        actions.append(["序号", "议题", "客户侧准备", "DCForge 侧准备", "共同输出", "状态"])
        for index, question in enumerate(self.brief.open_questions, 1):
            actions.append([
                index,
                question,
                "安排对应业务、IT 或安全负责人并准备相关资料",
                "准备访谈提纲、模板和现有能力说明",
                "写入双方确认的 RequirementBaseline 或项目附件",
                "待安排",
            ])
        configure(actions, [10, 52, 44, 42, 48, 16])
        for sheet in workbook.worksheets:
            sheet.sheet_view.showGridLines = False
            sheet.sheet_properties.pageSetUpPr.fitToPage = True
            sheet.page_setup.fitToWidth = 1
            sheet.page_setup.fitToHeight = 0
            sheet.sheet_properties.tabColor = blue
        workbook.save(path)

    def _html(self) -> str:
        def items(values: list[str]) -> str:
            return "".join(f"<li>{escape(value)}</li>" for value in values)

        layer_html = "".join(
            f"""<article class="layer"><h3>{escape(layer['layer'])}</h3>
<p class="components">{escape(' · '.join(layer['components']))}</p><p>{escape(layer['purpose'])}</p></article>"""
            for layer in self.brief.solution_layers
        )
        scenario_html = "".join(
            f"""<article class="scenario"><div><span class="eyebrow">重点场景</span><h3>{escape(s.name)}</h3>
<p><b>当前做法：</b>{escape(s.current_state)}</p></div><div><h4>AI 可以承担</h4><ul>{items(s.ai_assistance)}</ul></div>
<div class="boundary"><h4>人工边界</h4><p>{escape(s.human_boundary)}</p></div>
<div><h4>交付证据</h4><p>{escape('、'.join(s.evidence_output))}</p></div></article>"""
            for s in self.brief.scenarios
        )
        phases = "".join(
            f"""<article class="phase"><span>{escape(p.phase)}</span><h3>{escape(p.objective)}</h3>
<ul>{items(p.work)}</ul><p><b>输出：</b>{escape('、'.join(p.outputs))}</p>
<p class="meta">{escape(p.recommended_duration)}<br>{escape(p.prerequisite)}</p></article>"""
            for p in self.brief.implementation_phases
        )
        metrics = "".join(
            f"<tr><td>{escape(m.dimension)}</td><td>{escape(m.metric)}</td><td>{escape(m.proposed_threshold)}</td><td>{escape(m.evidence)}</td></tr>"
            for m in self.brief.acceptance_metrics
        )
        options = "".join(
            f"""<article class="option"><h3>{escape(o.name)}</h3><p>{escape(o.fit)}</p><ul>{items(o.scope)}</ul>
<p><b>前提：</b>{escape(o.prerequisite)}</p><p class="recommend">{escape(o.recommendation)}</p></article>"""
            for o in self.brief.package_options
        )
        return f"""<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{escape(self.brief.solution_title)}</title>
<style>
:root{{--navy:#0b1739;--blue:#1e5eff;--cyan:#18b7c9;--green:#10b981;--orange:#f59e0b;--ink:#172033;--muted:#5d6c85;--line:#dbe3f0;--light:#f4f7fb;font-family:Inter,"PingFang SC","Microsoft YaHei",sans-serif;color:var(--ink);background:var(--light)}}
*{{box-sizing:border-box}}body{{margin:0}}header{{min-height:88vh;background:linear-gradient(140deg,#081329,#102658 65%,#114f67);color:#fff;display:flex;align-items:center;padding:7vw}}header .inner{{max-width:1060px}}.brand{{color:var(--cyan);letter-spacing:.12em;font-weight:800}}h1{{font-size:clamp(36px,6vw,72px);line-height:1.08;max-width:1000px;margin:.6em 0 .35em}}header p{{font-size:20px;line-height:1.7;color:#d5e1f7;max-width:850px}}.status{{display:inline-block;background:#17366f;border:1px solid #31599f;border-radius:999px;padding:10px 18px;margin:20px 0}}main{{max-width:1180px;margin:auto;padding:64px 28px}}section{{margin:0 0 76px}}h2{{font-size:34px;color:var(--navy);margin:0 0 24px}}h3{{color:var(--navy)}}.lead{{font-size:18px;line-height:1.8;color:var(--muted);max-width:920px}}.notice{{background:#fff8e7;border:1px solid #f2c76e;border-radius:12px;padding:18px;color:#7b4a0b;line-height:1.7}}.grid{{display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:18px}}.card,.layer,.phase,.option,.scenario{{background:#fff;border:1px solid var(--line);border-radius:16px;padding:22px;box-shadow:0 10px 25px rgba(22,42,86,.05)}}.card{{border-left:5px solid var(--blue)}}.card p,.layer p,.scenario p,.phase p,.option p,li{{line-height:1.7}}.layers{{display:grid;gap:13px}}.layer{{display:grid;grid-template-columns:220px 1fr 1fr;align-items:center}}.layer h3{{margin:0}}.components{{color:var(--blue);font-weight:700}}.workflow{{display:grid;grid-template-columns:repeat(4,1fr);gap:14px;counter-reset:step}}.workflow div{{background:var(--navy);color:#fff;border-radius:14px;padding:20px;min-height:108px;counter-increment:step}}.workflow div:before{{content:counter(step,decimal-leading-zero);display:block;color:var(--cyan);font-weight:800;margin-bottom:14px}}.scenarios{{display:grid;gap:22px}}.scenario{{display:grid;grid-template-columns:1.15fr 1.35fr;gap:14px 26px;border-top:5px solid var(--cyan)}}.boundary{{background:#fff1f2;border-radius:10px;padding:4px 16px}}.eyebrow{{color:var(--blue);font-size:12px;font-weight:800;letter-spacing:.1em}}.phases,.options{{display:grid;grid-template-columns:repeat(4,minmax(0,1fr));gap:16px}}.options{{grid-template-columns:repeat(3,minmax(0,1fr))}}.phase span{{color:var(--blue);font-weight:800}}.meta{{font-size:13px;color:var(--muted)}}table{{width:100%;border-collapse:collapse;background:#fff;border-radius:14px;overflow:hidden}}th{{background:var(--navy);color:#fff;text-align:left}}th,td{{padding:14px;border:1px solid var(--line);vertical-align:top;line-height:1.5}}tr:nth-child(even){{background:#f8faff}}.recommend{{color:var(--green);font-weight:800}}footer{{background:var(--navy);color:#d5e1f7;padding:44px 7vw;line-height:1.7}}.toolbar{{position:fixed;right:20px;bottom:20px}}button{{border:0;border-radius:999px;background:var(--blue);color:#fff;padding:13px 20px;font:inherit;cursor:pointer;box-shadow:0 8px 24px rgba(30,94,255,.28)}}
@media(max-width:850px){{.grid,.scenario,.phases,.options{{grid-template-columns:1fr}}.workflow{{grid-template-columns:repeat(2,1fr)}}.layer{{grid-template-columns:1fr}}}}
@media print{{header{{min-height:auto;padding:35px}}main{{padding:30px}}section{{break-inside:avoid;margin-bottom:38px}}.toolbar{{display:none}}}}
</style></head><body>
<header><div class="inner"><div class="brand">DCFORGE</div><h1>{escape(self.brief.solution_title)}</h1>
<p>从客户沟通到需求基线、企业知识、方案评审和客户发布的完整售前闭环</p>
<div class="status">{escape(self.brief.document_status)} · v{escape(self.brief.package_version)}</div>
<p>面向：{escape(self.brief.client_label)}<br>编制：{escape(self.brief.prepared_by)} · {escape(self.brief.generated_date)}</p></div></header>
<main>
<section><h2>材料说明</h2><p class="lead">本建议书可用于客户正式汇报、技术交流、试点立项和验收口径讨论。项目启动后，双方应通过 RequirementBaseline 确认正式需求。</p>
<div class="notice">材料中的目标值是建议验收口径，不是当前客户已经取得的业务成果。正式范围、周期、价格和服务等级以双方后续确认文件为准。</div></section>
<section><h2>我们目前如何理解贵司</h2><div class="grid">{''.join(f'<article class="card"><p>{escape(v)}</p></article>' for v in self.brief.known_current_state)}</div></section>
<section><h2>本次建议聚焦的问题</h2><div class="grid">{''.join(f'<article class="card"><p>{escape(v)}</p></article>' for v in self.brief.pain_points)}</div></section>
<section><h2>建设原则</h2><ul>{items(self.brief.principles)}</ul></section>
<section><h2>方案总览</h2><div class="layers">{layer_html}</div></section>
<section><h2>完整业务闭环</h2><div class="workflow">{''.join(f'<div>{escape(v)}</div>' for v in self.brief.workflow)}</div></section>
<section><h2>重点业务场景</h2><div class="scenarios">{scenario_html}</div></section>
<section><h2>数据、安全和人工决策边界</h2><div class="grid">{''.join(f'<article class="card"><p>{escape(v)}</p></article>' for v in self.brief.security_boundaries)}</div></section>
<section><h2>分阶段实施路线</h2><div class="phases">{phases}</div></section>
<section><h2>建议验收口径</h2><table><thead><tr><th>维度</th><th>指标</th><th>建议口径</th><th>验收证据</th></tr></thead><tbody>{metrics}</tbody></table>
<p class="notice">所有阈值均待双方确认。验收前应冻结样本集、计算口径、人工金标准和异常处理规则。</p></section>
<section><h2>三档建设范围</h2><div class="options">{options}</div></section>
<section><h2>双方待确认事项</h2><ul>{items(self.brief.open_questions)}</ul></section>
</main><footer><b>内容边界</b><ul>{items(self.brief.content_boundary)}</ul></footer>
<div class="toolbar"><button onclick="window.print()">打印或导出 PDF</button></div></body></html>"""

    def _security_markdown(self) -> str:
        rows = "\n".join(
            f"| {item['system']} | {item['read']} | {item['write']} | {item['mode']} | {item['boundary']} |"
            for item in self.brief.integration_objects
        )
        boundaries = "\n".join(f"- {value}" for value in self.brief.security_boundaries)
        return f"""# DCForge 部署、安全与集成边界

文档状态：{self.brief.document_status}  
版本：{self.brief.package_version}  
日期：{self.brief.generated_date}

## 使用说明

本文用于部署架构、安全评审和系统接口讨论。具体配置应在客户安全、IT、采购和合规负责人确认后写入项目实施附件。

## 安全边界

{boundaries}

## 系统集成建议

| 系统 | 建议读取 | 建议回写 | 接入方式 | 业务与安全边界 |
|---|---|---|---|---|
{rows}

## 身份和权限

1. 客户用户只能访问本人被授权的项目和正式发布内容。
2. 企业员工按项目、岗位和知识资产权限访问内部工作台。
3. Agent 调用知识库、MCP 和业务工具时携带用户、项目、时间和权限上下文。
4. 未获得来源权限时，系统返回权限拒绝或证据不足，不使用其他项目资料补答。

## 数据处理

1. 采购材料、报价、合同和供应商信息在客户授权的数据域内处理。
2. 敏感字段按岗位、场景和输出渠道脱敏。
3. 模型输入、知识检索、工具调用和输出保留必要审计信息。
4. 数据保留期限、日志字段、备份和删除机制由客户安全要求确定。

## 人工决策门禁

AI 可以做材料检查、信息提取、规则匹配、风险提示和草稿生成。供应商准入、采购审批、定标、合同签署和付款仍由企业人员及原有系统完成。任何需要人工审批的事项都应保留审批人、时间、意见和版本记录。

## 上线前检查

- 完成部署区域、网络和模型访问方式评审。
- 完成身份源、单点登录、岗位权限和项目权限配置。
- 完成 SRM、ERP、OA 和飞书接口授权。
- 完成数据分类、脱敏、日志、备份和删除策略确认。
- 完成权限负向测试、人工门禁测试和审计记录抽查。
- 完成客户安全、采购、合规和 IT 联合签字。
"""

    def _demo_markdown(self) -> str:
        return f"""# DCForge 演示与操作手册

文档状态：{self.brief.document_status}  
版本：{self.brief.package_version}

## 演示目标

用一个大型汽车制造企业采购场景说明完整闭环：客户在飞书描述问题，Agent 形成需求状态，企业员工补充知识与情报，系统生成方案，内部评审后发布给客户，客户反馈再进入下一版需求。

## 演示角色

- 客户代表：在飞书沟通需求，并在客户需求与方案中心确认或补充信息。
- 售前负责人：在统一售前工作台查看会话、状态、缺口、资料和方案。
- 方案评审人：批准或驳回当前成果修订。
- 系统管理员：准备演示账号、项目权限和知识资料。

## 演示前准备

1. 确认 API 和飞书长连接均处于可用状态。
2. 由项目经理提供正式工作台地址和客户访问链接。
3. 准备一份脱敏供应商材料、一份采购制度和一份外部政策资料。
4. 确认演示项目没有学校、具体车型或采购数量等无关测试输入。
5. 确认演示指标标记为建议值或待验证值。

## 标准演示脚本

### 1. 客户在飞书描述问题

建议发送：

> 我们是一家大型汽车制造企业，目前供应商准入、询比价和采购合规审查主要依赖人工。采购周期较长，流程留痕也不完整，希望了解 AI 可以在哪些环节提供帮助。

讲解重点：机器人面向客户只做自然回复，不显示内部状态版本、成熟度、技能名和后台字段。

### 2. 企业员工打开统一售前工作台

展示客户会话、RequirementState、缺口、冲突和 Baseline 状态。说明飞书项目会自动进入工作台，不需要重复建档。

### 3. 录入资料和外部情报

录入客户材料时，内容进入 Requirement Intelligence。录入企业资料或外部情报时，内容只作为研究证据，不自动变成客户确认需求。

### 4. 运行知识与情报研究

展示企业知识结果、外部来源 URL 和引用。说明知识检索同时检查项目、用户、时间和权限。

### 5. 形成需求基线

客户或授权员工确认需求后形成 RequirementBaseline。没有正式 Baseline 时，系统拒绝生成正式方案草稿。

### 6. 生成和编辑方案

生成三档方案和客户成果稿。编辑成果后修订号增加，之前的内部批准自动失效。

### 7. 内部评审和客户发布

评审人批准当前修订后，售前负责人发布方案。客户页面只显示正式需求、正式方案和客户安全成果。

### 8. 客户反馈迭代

客户在需求与方案中心补充或纠正信息。反馈进入 Requirement Intelligence，形成新的 State 版本并显示差异。

## 演示失败处理

- 飞书无回复：检查长连接、应用权限、模型配置和机器人允许用户范围。
- 工作台看不到项目：检查客户消息是否已经持久化，以及内部访问令牌是否正确。
- 无法生成方案：先检查正式 Baseline 和研究快照是否存在。
- 无法发布：检查当前成果修订是否通过内部批准。
- 客户打不开成果：重新获取客户访问链接，并确认访问令牌没有过期。

## 演示结束口径

本次演示证明系统链路和治理边界可以运行，不代表客户业务指标已经达成。试点验收需要双方确认样本、基线、计算口径和目标值。
"""

    def _readme_markdown(self) -> str:
        file_rows = [
            (PPTX_NAME, "客户汇报主材料，可直接编辑"),
            (PDF_NAME, "适合邮件发送、归档和打印"),
            (HTML_NAME, "离线可浏览，可打印或另存为 PDF"),
            (XLSX_NAME, "实施计划、验收指标、接口清单和双方待办"),
            (SECURITY_NAME, "供客户 IT、数据安全和架构团队评审"),
            (DEMO_NAME, "供售前、交付和客户演示人员使用"),
            (MANIFEST_NAME, "文件大小、SHA-256 和内容边界"),
        ]
        rows = "\n".join(f"| `{name}` | {description} |" for name, description in file_rows)
        boundary = "\n".join(f"- {value}" for value in self.brief.content_boundary)
        return f"""# DCForge 企业智能招采解决方案交付包

版本：{self.brief.package_version}  
文档状态：{self.brief.document_status}  
编制日期：{self.brief.generated_date}

## 交付对象

本包面向大型汽车制造企业的采购、合规、法务、IT、数据安全和项目决策人员，可直接用于首次正式汇报、技术交流、试点立项和验收口径讨论。

## 文件说明

| 文件 | 用途 |
|---|---|
{rows}

## 建议阅读顺序

1. 决策人与业务负责人先看 PPT 或 PDF。
2. IT 和安全负责人查看部署安全与集成边界。
3. 项目经理打开实施与验收计划，确认负责人、资料和排期。
4. 演示人员按操作手册准备飞书、工作台和客户中心演示。

## 内容边界

{boundary}

## 完整性校验

`manifest.json` 记录每个文件的 SHA-256。交付前可用企业文件校验工具核对，确认文件没有在传输过程中被修改。
"""


def build_default_enterprise_delivery(output_dir: Path | str) -> DeliveryBuildResult:
    return EnterpriseDeliveryPackageBuilder(
        EnterpriseDeliveryBrief.default_automotive_procurement()
    ).build(output_dir)


if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(description="Build the DCForge enterprise handoff package")
    parser.add_argument("output_dir", nargs="?", default="output/enterprise-delivery-v1.0")
    args = parser.parse_args()
    result = build_default_enterprise_delivery(args.output_dir)
    print(result.output_dir)
    print(result.archive_path)
